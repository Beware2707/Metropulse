"""A tiny scene graph, so the deck is described once and rendered twice.

There is no LibreOffice in this environment, so the PDF of the deck cannot be
produced by converting the PPTX. Writing the deck twice would guarantee the two
drift apart — and a submission where the slides and the PDF disagree is worse
than either alone. So the deck emits drawing operations, and two backends
render them: python-pptx for the editable file, reportlab for the PDF.
"""

from __future__ import annotations

import os
from dataclasses import dataclass, field

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from reportlab.lib.colors import HexColor
from reportlab.lib.utils import ImageReader
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfgen import canvas as pdfcanvas

from . import theme as T

W, H = 13.333, 7.5
LEFT, CENTER, RIGHT = "l", "c", "r"
TOP, MIDDLE = "t", "m"


@dataclass
class Op:
    kind: str
    args: dict


@dataclass
class Slide:
    ops: list[Op] = field(default_factory=list)


class Recorder:
    """Collects drawing operations. Geometry is in inches, origin top-left."""

    def __init__(self) -> None:
        self.slides: list[Slide] = []
        self._cur: Slide | None = None

    # ------------------------------------------------------------ surface
    def slide(self, dark: bool = False):
        self._cur = Slide()
        self.slides.append(self._cur)
        self.rect(0, 0, W, H, T.INK if dark else T.PAPER)
        return self._cur

    def _add(self, kind: str, **args) -> None:
        assert self._cur is not None, "call slide() first"
        self._cur.ops.append(Op(kind, args))

    # --------------------------------------------------------- primitives
    def rect(self, x, y, w, h, fill, radius=False):
        self._add("rect", x=x, y=y, w=w, h=h, fill=fill, radius=radius)

    def text(self, x, y, w, h, body, *, size=12, bold=False, colour=T.BODY,
             font=None, align=LEFT, anchor=TOP, spacing=1.0):
        self._add("text", x=x, y=y, w=w, h=h, body=body, size=size, bold=bold,
                  colour=colour, font=font or T.BODYF, align=align,
                  anchor=anchor, spacing=spacing)

    def bullets(self, x, y, w, h, items, *, size=11.5, colour=T.BODY, gap=7):
        self._add("bullets", x=x, y=y, w=w, h=h, items=items, size=size,
                  colour=colour, gap=gap)

    def pill(self, x, y, label):
        w = 0.30 + len(label) * 0.077
        self._add("pill", x=x, y=y, w=w, label=label,
                  colour=T.STATUS.get(label, T.SLATE))
        return w

    def arrow(self, x, y, w, h=0.16):
        self._add("arrow", x=x, y=y, w=w, h=h)

    def image(self, path, x, y, *, h):
        self._add("image", path=path, x=x, y=y, h=h)

    # ----------------------------------------------------------- compound
    def step(self, x, y, w, h, title, sub, colour, light=True):
        """One box in a flow diagram.

        `light` describes the *text*, not the box: True paints white text for a
        dark or saturated fill, False paints navy text for a pale one. Passing
        it the wrong way round renders navy on navy, which disappears rather
        than merely looking wrong.
        """
        self.rect(x, y, w, h, colour, radius=True)
        self.text(x + 0.1, y + 0.13, w - 0.2, 0.46, title, size=11.5, bold=True,
                  font=T.HEAD, colour="FFFFFF" if light else T.INK, align=CENTER)
        if sub:
            self.text(x + 0.09, y + 0.6, w - 0.18, h - 0.68, sub, size=9,
                      colour="D7E0F2" if light else T.BODY, align=CENTER)

    def header(self, kicker, title, num):
        self.text(LEFTM := 0.7, 0.42, 9.0, 0.26, kicker.upper(), size=10.5,
                  bold=True, colour=T.PRIMARY, font=T.HEAD)
        self.text(LEFTM, 0.72, W - 2.3, 0.76, title, size=29, bold=True,
                  colour=T.INK, font=T.HEAD)
        self.text(W - 1.4, 0.42, 0.7, 0.28, num, size=11, colour=T.MUTED,
                  font=T.HEAD, align=RIGHT)

    def foot(self, txt):
        self.text(0.7, H - 0.66, W - 1.4, 0.36, txt, size=9.5, colour=T.MUTED)


# =========================================================== pptx renderer
def _rgb(hexstr: str) -> RGBColor:
    return RGBColor(*T.rgb(hexstr))


_ALIGN_P = {LEFT: PP_ALIGN.LEFT, CENTER: PP_ALIGN.CENTER, RIGHT: PP_ALIGN.RIGHT}
_ANCHOR_P = {TOP: MSO_ANCHOR.TOP, MIDDLE: MSO_ANCHOR.MIDDLE}


def render_pptx(rec: Recorder, out_path: str) -> str:
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    blank = prs.slide_layouts[6]

    for sl in rec.slides:
        s = prs.slides.add_slide(blank)
        for op in sl.ops:
            a = op.args
            if op.kind in ("rect", "pill"):
                is_pill = op.kind == "pill"
                shp = s.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE if (a.get("radius") or is_pill)
                    else MSO_SHAPE.RECTANGLE,
                    Inches(a["x"]), Inches(a["y"]), Inches(a["w"]),
                    Inches(0.27 if is_pill else a["h"]))
                shp.fill.solid()
                shp.fill.fore_color.rgb = _rgb(a["colour"] if is_pill else a["fill"])
                shp.line.fill.background()
                shp.shadow.inherit = False
                try:
                    shp.adjustments[0] = 0.5 if is_pill else 0.08
                except (IndexError, AttributeError):
                    pass
                if is_pill:
                    tf = shp.text_frame
                    tf.margin_left = tf.margin_right = 0
                    tf.margin_top = tf.margin_bottom = 0
                    tf.word_wrap = False
                    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                    p = tf.paragraphs[0]
                    p.alignment = PP_ALIGN.CENTER
                    r = p.add_run()
                    r.text = a["label"]
                    r.font.size = Pt(8)
                    r.font.bold = True
                    r.font.name = T.HEAD
                    r.font.color.rgb = _rgb("FFFFFF")
            elif op.kind == "arrow":
                shp = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(a["x"]),
                                         Inches(a["y"]), Inches(a["w"]), Inches(a["h"]))
                shp.fill.solid()
                shp.fill.fore_color.rgb = _rgb("AFC0DC")
                shp.line.fill.background()
                shp.shadow.inherit = False
            elif op.kind == "text":
                box = s.shapes.add_textbox(Inches(a["x"]), Inches(a["y"]),
                                           Inches(a["w"]), Inches(a["h"]))
                tf = box.text_frame
                tf.word_wrap = True
                tf.margin_left = tf.margin_right = 0
                tf.margin_top = tf.margin_bottom = 0
                tf.vertical_anchor = _ANCHOR_P[a["anchor"]]
                for i, line in enumerate(str(a["body"]).split("\n")):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.alignment = _ALIGN_P[a["align"]]
                    if a["spacing"] != 1.0:
                        p.line_spacing = a["spacing"]
                    r = p.add_run()
                    r.text = line
                    r.font.size = Pt(a["size"])
                    r.font.bold = a["bold"]
                    r.font.name = a["font"]
                    r.font.color.rgb = _rgb(a["colour"])
            elif op.kind == "bullets":
                box = s.shapes.add_textbox(Inches(a["x"]), Inches(a["y"]),
                                           Inches(a["w"]), Inches(a["h"]))
                tf = box.text_frame
                tf.word_wrap = True
                tf.margin_left = tf.margin_right = 0
                tf.margin_top = tf.margin_bottom = 0
                for i, item in enumerate(a["items"]):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.space_after = Pt(a["gap"])
                    r = p.add_run()
                    r.text = "•  " + item
                    r.font.size = Pt(a["size"])
                    r.font.name = T.BODYF
                    r.font.color.rgb = _rgb(a["colour"])
            elif op.kind == "image":
                s.shapes.add_picture(a["path"], Inches(a["x"]), Inches(a["y"]),
                                     height=Inches(a["h"]))
    prs.save(out_path)
    return out_path


# ============================================================ pdf renderer
PT = 72.0
_FONT = {(T.HEAD, False): "Helvetica", (T.HEAD, True): "Helvetica-Bold",
         (T.BODYF, False): "Helvetica", (T.BODYF, True): "Helvetica-Bold"}


def _wrap(text: str, font: str, size: float, width_pt: float) -> list[str]:
    out: list[str] = []
    for para in str(text).split("\n"):
        words, line = para.split(), ""
        if not words:
            out.append("")
            continue
        for wd in words:
            trial = wd if not line else line + " " + wd
            if pdfmetrics.stringWidth(trial, font, size) <= width_pt:
                line = trial
            else:
                if line:
                    out.append(line)
                line = wd
        out.append(line)
    return out


def render_pdf(rec: Recorder, out_path: str) -> str:
    c = pdfcanvas.Canvas(out_path, pagesize=(W * PT, H * PT))
    for sl in rec.slides:
        for op in sl.ops:
            a = op.args
            X, Y = a.get("x", 0) * PT, a.get("y", 0) * PT

            def top(y_in):                      # inches from top -> pdf y
                return H * PT - y_in * PT

            if op.kind in ("rect", "pill"):
                is_pill = op.kind == "pill"
                h = 0.27 if is_pill else a["h"]
                c.setFillColor(HexColor("#" + (a["colour"] if is_pill else a["fill"])))
                if a.get("radius") or is_pill:
                    c.roundRect(X, top(a["y"] + h), a["w"] * PT, h * PT,
                                (h * PT / 2) if is_pill else 6, fill=1, stroke=0)
                else:
                    c.rect(X, top(a["y"] + h), a["w"] * PT, h * PT, fill=1, stroke=0)
                if is_pill:
                    c.setFillColor(HexColor("#FFFFFF"))
                    c.setFont("Helvetica-Bold", 8)
                    c.drawCentredString(X + a["w"] * PT / 2,
                                        top(a["y"] + h) + h * PT / 2 - 2.8, a["label"])
            elif op.kind == "arrow":
                c.setFillColor(HexColor("#AFC0DC"))
                c.rect(X, top(a["y"] + a["h"]), a["w"] * PT, a["h"] * PT,
                       fill=1, stroke=0)
            elif op.kind == "text":
                font = _FONT[(a["font"], a["bold"])]
                size = a["size"]
                c.setFont(font, size)
                c.setFillColor(HexColor("#" + a["colour"]))
                lines = _wrap(a["body"], font, size, a["w"] * PT)
                lh = size * 1.22 * a["spacing"]
                block = len(lines) * lh
                y0 = (top(a["y"]) - size
                      if a["anchor"] == TOP
                      else top(a["y"] + a["h"] / 2) + block / 2 - size)
                for i, line in enumerate(lines):
                    yy = y0 - i * lh
                    if a["align"] == CENTER:
                        c.drawCentredString(X + a["w"] * PT / 2, yy, line)
                    elif a["align"] == RIGHT:
                        c.drawRightString(X + a["w"] * PT, yy, line)
                    else:
                        c.drawString(X, yy, line)
            elif op.kind == "bullets":
                size = a["size"]
                c.setFillColor(HexColor("#" + a["colour"]))
                yy = top(a["y"]) - size
                for item in a["items"]:
                    lines = _wrap(item, "Helvetica", size, a["w"] * PT - 12)
                    for j, line in enumerate(lines):
                        c.setFont("Helvetica", size)
                        if j == 0:
                            c.drawString(X, yy, "•")
                        c.drawString(X + 12, yy, line)
                        yy -= size * 1.24
                    yy -= a["gap"] * 0.75
            elif op.kind == "image":
                img = ImageReader(a["path"])
                iw, ih = img.getSize()
                w_in = a["h"] * (iw / ih)
                c.drawImage(img, X, top(a["y"] + a["h"]), w_in * PT, a["h"] * PT,
                            mask="auto")
        c.showPage()
    c.save()
    return out_path
